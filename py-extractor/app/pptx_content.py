from __future__ import annotations

import io
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app import markdown_table


@dataclass
class PptxContent:
    text: str
    slide_count: int
    table_count: int
    warnings: list[str]


def extract(data: bytes) -> PptxContent:
    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    table_count = 0
    warnings: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            shape_blocks, shape_tables, shape_warnings = _shape_content(shape)
            blocks.extend(shape_blocks)
            table_count += shape_tables
            for warning in shape_warnings:
                if warning not in warnings:
                    warnings.append(warning)
        if blocks:
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(blocks))
    return PptxContent(
        text="\n\n".join(slides),
        slide_count=len(presentation.slides),
        table_count=table_count,
        warnings=warnings,
    )


def _shape_content(shape: object) -> tuple[list[str], int, list[str]]:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        blocks: list[str] = []
        tables = 0
        warnings: list[str] = []
        for child in shape.shapes:
            child_blocks, child_tables, child_warnings = _shape_content(child)
            blocks.extend(child_blocks)
            tables += child_tables
            warnings.extend(child_warnings)
        return blocks, tables, warnings

    if getattr(shape, "has_table", False):
        rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
        table = markdown_table.serialize(rows)
        return ([table] if table else []), (1 if table else 0), []

    text = str(getattr(shape, "text", "") or "").strip()
    if text:
        return [text], 0, []

    warning_by_type = {
        MSO_SHAPE_TYPE.CHART: "pptx_chart_not_extracted",
        MSO_SHAPE_TYPE.DIAGRAM: "pptx_diagram_not_extracted",
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "pptx_ole_not_extracted",
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "pptx_ole_not_extracted",
    }
    warning = warning_by_type.get(shape_type)
    return [], 0, [warning] if warning else []

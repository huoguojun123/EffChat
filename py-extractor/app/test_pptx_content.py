from __future__ import annotations

import io
import unittest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from app import main, pptx_content


class PptxContentTest(unittest.TestCase):
    def test_table_group_order_and_multiple_slides(self):
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[6])
        first.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "before"
        table = first.shapes.add_table(2, 2, Inches(0), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "H|1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "V1"
        table.cell(1, 1).text = "V2"
        group = first.shapes.add_group_shape()
        group.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "group-one"
        nested = group.shapes.add_group_shape()
        nested.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "group-two"

        second = presentation.slides.add_slide(presentation.slide_layouts[6])
        second.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "second-slide"

        buffer = io.BytesIO()
        presentation.save(buffer)
        extracted = main.extract_pptx(buffer.getvalue(), "fixture.pptx")

        self.assertEqual(extracted.page_count, 2)
        self.assertEqual(extracted.table_count, 1)
        self.assertEqual(extracted.warnings, [])
        self.assertLess(extracted.text.index("before"), extracted.text.index(r"H\|1"))
        self.assertLess(extracted.text.index(r"H\|1"), extracted.text.index("group-one"))
        self.assertLess(extracted.text.index("group-one"), extracted.text.index("group-two"))
        self.assertLess(extracted.text.index("group-two"), extracted.text.index("## Slide 2"))
        self.assertIn("second-slide", extracted.text)

    def test_table_only_slide_and_bounded_unsupported_warning(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(2, 1, Inches(0), Inches(0), Inches(2), Inches(2)).table
        table.cell(0, 0).text = "Header"
        table.cell(1, 0).text = "Value"
        buffer = io.BytesIO()
        presentation.save(buffer)

        extracted = main.extract_pptx(buffer.getvalue(), "table-only.pptx")
        self.assertIn("## Slide 1", extracted.text)
        self.assertIn("| Header |", extracted.text)
        self.assertEqual(extracted.table_count, 1)

        class UnsupportedShape:
            shape_type = MSO_SHAPE_TYPE.CHART
            has_table = False
            text = ""

        blocks, tables, warnings = pptx_content._shape_content(UnsupportedShape())
        self.assertEqual((blocks, tables, warnings), ([], 0, ["pptx_chart_not_extracted"]))


if __name__ == "__main__":
    unittest.main()

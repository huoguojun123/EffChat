from __future__ import annotations

import asyncio
import csv
import io
import threading
import time
import unittest
import zipfile
from unittest import mock

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app import main, resource_limits


def office_zip(entries: list[tuple[str, bytes]], compression: int = zipfile.ZIP_STORED) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=compression) as archive:
        for name, content in entries:
            archive.writestr(name, content)
    return buffer.getvalue()


class FakeUpload:
    filename = "fixture.csv"
    content_type = "text/csv"

    def __init__(self, data: bytes):
        self.data = data

    async def read(self, size: int = -1) -> bytes:
        return self.data if size < 0 else self.data[:size]


class ExtractorResourceLimitsTest(unittest.TestCase):
    def test_office_archive_rejects_each_resource_limit(self):
        with (
            mock.patch.object(resource_limits, "OFFICE_MAX_ENTRIES", 2),
            mock.patch.object(resource_limits, "OFFICE_MAX_ENTRY_BYTES", 2),
            mock.patch.object(resource_limits, "OFFICE_MAX_UNCOMPRESSED_BYTES", 4),
        ):
            resource_limits.validate_office_archive(office_zip([("a", b"12"), ("b", b"34")]))

        too_many_entries = office_zip([("a", b"1"), ("b", b"2")])
        with mock.patch.object(resource_limits, "OFFICE_MAX_ENTRIES", 1):
            with mock.patch.object(resource_limits.zipfile, "ZipFile") as zip_reader:
                with self.assertRaisesRegex(main.HTTPException, "office_archive_entry_limit_exceeded"):
                    resource_limits.validate_office_archive(too_many_entries)
                zip_reader.assert_not_called()

        with mock.patch.object(resource_limits, "OFFICE_MAX_CENTRAL_DIRECTORY_BYTES", 1):
            with self.assertRaisesRegex(main.HTTPException, "office_archive_directory_limit_exceeded"):
                resource_limits.validate_office_archive(office_zip([("a", b"1")]))

        with mock.patch.object(resource_limits, "OFFICE_MAX_ENTRY_BYTES", 3):
            with self.assertRaisesRegex(main.HTTPException, "office_archive_file_limit_exceeded"):
                resource_limits.validate_office_archive(office_zip([("a", b"1234")]))

        with mock.patch.object(resource_limits, "OFFICE_MAX_UNCOMPRESSED_BYTES", 3):
            with self.assertRaisesRegex(main.HTTPException, "office_archive_size_limit_exceeded"):
                resource_limits.validate_office_archive(office_zip([("a", b"12"), ("b", b"34")]))

        compressed = office_zip([("a", b"a" * 4096)], compression=zipfile.ZIP_DEFLATED)
        with mock.patch.object(resource_limits, "OFFICE_MAX_COMPRESSION_RATIO", 2):
            with self.assertRaisesRegex(main.HTTPException, "office_archive_ratio_limit_exceeded"):
                resource_limits.validate_office_archive(compressed)

    def test_normal_office_fixtures_still_extract(self):
        docx_buffer = io.BytesIO()
        docx = Document()
        docx.add_paragraph("DOCX fixture")
        docx.save(docx_buffer)
        self.assertIn("DOCX fixture", main.extract_docx(docx_buffer.getvalue(), "fixture.docx").text)

        pptx_buffer = io.BytesIO()
        pptx = Presentation()
        slide = pptx.slides.add_slide(pptx.slide_layouts[5])
        slide.shapes.title.text = "PPTX fixture"
        pptx.save(pptx_buffer)
        self.assertIn("PPTX fixture", main.extract_pptx(pptx_buffer.getvalue(), "fixture.pptx").text)

        xlsx_buffer = io.BytesIO()
        workbook = Workbook()
        workbook.active.append(["XLSX fixture", "value"])
        workbook.save(xlsx_buffer)
        workbook.close()
        self.assertIn("XLSX fixture", main.extract_xlsx(xlsx_buffer.getvalue(), "fixture.xlsx").text)

    def test_csv_preserves_long_single_column_content(self):
        value = "a" * 140_000
        extracted = main.extract_csv(value.encode(), "fixture.csv")
        self.assertIn(value, extracted.text)
        self.assertEqual(extracted.table_count, 1)

    def test_csv_rejects_explicit_limits_and_uncertain_delimiter(self):
        with mock.patch.object(resource_limits, "CSV_MAX_ROWS", 1):
            with self.assertRaisesRegex(main.HTTPException, "csv_row_limit_exceeded"):
                main.extract_csv(b"a\nb\n", "rows.csv")

        with mock.patch.object(resource_limits, "CSV_MAX_COLUMNS", 2):
            with self.assertRaisesRegex(main.HTTPException, "csv_column_limit_exceeded"):
                main.extract_csv(b"a,b,c\n", "wide.csv")

        with mock.patch.object(resource_limits, "CSV_MAX_CELLS", 3):
            with self.assertRaisesRegex(main.HTTPException, "csv_cell_limit_exceeded"):
                main.extract_csv(b"a,b\nc,d\n", "cells.csv")

        with mock.patch.object(main, "MAX_OUTPUT_BYTES", 3):
            with self.assertRaisesRegex(main.HTTPException, "csv_content_limit_exceeded"):
                main.extract_csv(b"abcd\n", "content.csv")

        previous_limit = csv.field_size_limit(resource_limits.CSV_MAX_FIELD_CHARS)
        try:
            oversized = ("x" * (resource_limits.CSV_MAX_FIELD_CHARS + 1)).encode()
            with self.assertRaisesRegex(main.HTTPException, "csv_field_limit_exceeded"):
                main.extract_csv(oversized, "field.csv")
        finally:
            csv.field_size_limit(previous_limit)

        with mock.patch.object(resource_limits.csv.Sniffer, "sniff", side_effect=csv.Error("ambiguous")):
            with self.assertRaisesRegex(main.HTTPException, "csv_delimiter_uncertain"):
                main.extract_csv(b"a,b;c\n1,2;3\n", "ambiguous.csv")


class ExtractorConcurrencyTest(unittest.IsolatedAsyncioTestCase):
    async def test_slow_parsers_do_not_block_health_or_each_other(self):
        lock = threading.Lock()
        both_started = threading.Event()
        release = threading.Event()
        started_count = 0

        def slow_parser(data: bytes, filename: str) -> main.ExtractedDocument:
            nonlocal started_count
            with lock:
                started_count += 1
                if started_count == 2:
                    both_started.set()
            if not both_started.wait(timeout=0.5):
                raise RuntimeError("second parser never started")
            if not release.wait(timeout=0.5):
                raise RuntimeError("parser was not released")
            return main.ExtractedDocument(text="ok", parser="slow")

        with (
            mock.patch.object(main, "parser_for", return_value=slow_parser),
            mock.patch.object(main, "LOCAL_PARSE_SLOTS", asyncio.Semaphore(2)),
            mock.patch.object(main, "LOCAL_PARSE_QUEUE_TIMEOUT_SECONDS", 0.02),
        ):
            started_at = time.perf_counter()
            first = asyncio.create_task(main.extract(FakeUpload(b"one"), filename="", content_type=""))
            second = asyncio.create_task(main.extract(FakeUpload(b"two"), filename="", content_type=""))
            try:
                for _ in range(20):
                    if both_started.is_set():
                        break
                    await asyncio.sleep(0.005)
                self.assertTrue(both_started.is_set())
                self.assertEqual(await main.health(), {"status": "ok"})
                self.assertLess(time.perf_counter() - started_at, 0.2)
                with self.assertRaisesRegex(main.HTTPException, "extractor_busy") as busy:
                    await main.extract(FakeUpload(b"three"), filename="", content_type="")
                self.assertEqual(busy.exception.status_code, 503)
            finally:
                release.set()
            first_result, second_result = await asyncio.gather(first, second)
            recovered = await main.extract(FakeUpload(b"four"), filename="", content_type="")

        self.assertEqual(first_result["text"], "ok")
        self.assertEqual(second_result["text"], "ok")
        self.assertEqual(recovered["text"], "ok")
        self.assertTrue(both_started.is_set())


if __name__ == "__main__":
    unittest.main()
